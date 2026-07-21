"""Unit tests for file_extract — verifies text extraction across all
supported formats without hitting the API. Run with:
    cd /app/backend && pytest tests/test_file_extract.py -v
"""
import io
import pytest

import file_extract


def _make_docx(text: str) -> bytes:
    from docx import Document
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(rows: list[list]) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slide_titles: list[str]) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    for title in slide_titles:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
        slide.shapes.title.text = title
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pdf(text: str) -> bytes:
    # pypdf can't create PDFs — write via reportlab if available, else skip.
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    for i, line in enumerate(text.split("\n"), start=1):
        c.drawString(72, 800 - 15 * i, line)
    c.save()
    return buf.getvalue()


def test_csv_extraction():
    data = b"name,email\nAlice,alice@x.io\nBob,bob@x.io"
    text, truncated, mime = file_extract.extract_text("contacts.csv", data)
    assert "Alice,alice@x.io" in text
    assert "Bob,bob@x.io" in text
    assert not truncated
    assert mime == "text/csv"


def test_csv_utf8_bom():
    data = "\ufeffname,city\nRobin,Amsterdam".encode("utf-8")
    text, _, _ = file_extract.extract_text("bom.csv", data)
    assert "Robin,Amsterdam" in text
    # BOM stripped
    assert "\ufeff" not in text


def test_docx_extraction():
    data = _make_docx("Quarterly review\nRevenue up 20%")
    text, _, mime = file_extract.extract_text("q4.docx", data)
    assert "Quarterly review" in text
    assert "Revenue up 20%" in text
    assert mime.endswith("wordprocessingml.document")


def test_xlsx_extraction():
    data = _make_xlsx([["Product", "Price"], ["Widget", 9.99], ["Gadget", 19.99]])
    text, _, mime = file_extract.extract_text("prices.xlsx", data)
    assert "Product,Price" in text
    assert "Widget,9.99" in text
    assert "Sheet: Sheet1" in text
    assert mime.endswith("spreadsheetml.sheet")


def test_pptx_extraction():
    data = _make_pptx(["Kickstart Plan", "Roadmap Q1"])
    text, _, mime = file_extract.extract_text("deck.pptx", data)
    assert "Kickstart Plan" in text
    assert "Roadmap Q1" in text
    assert "Slide 1" in text
    assert mime.endswith("presentationml.presentation")


def test_pdf_extraction():
    reportlab = pytest.importorskip("reportlab")  # noqa: F841
    data = _make_pdf("Hello Zynthoro\nSecond line")
    text, _, mime = file_extract.extract_text("hello.pdf", data)
    assert "Hello Zynthoro" in text
    assert mime == "application/pdf"


def test_unsupported_extension_raises():
    with pytest.raises(ValueError, match="Unsupported"):
        file_extract.extract_text("virus.exe", b"MZ...")


def test_no_extension_raises():
    with pytest.raises(ValueError, match="Unsupported"):
        file_extract.extract_text("README", b"hi")


def test_corrupted_docx_raises_runtime():
    with pytest.raises(RuntimeError):
        file_extract.extract_text("broken.docx", b"not-a-real-docx")


def test_truncation():
    # Build a CSV bigger than MAX_CHARS
    big = ("a" * 500 + "\n") * ((file_extract.MAX_CHARS // 500) + 5)
    text, truncated, _ = file_extract.extract_text("big.csv", big.encode())
    assert truncated
    assert "truncated by Zynthoro" in text
    assert len(text) <= file_extract.MAX_CHARS + 200  # trailing note buffer


def test_case_insensitive_extension():
    data = b"a,b\n1,2"
    text, _, _ = file_extract.extract_text("Data.CSV", data)
    assert "1,2" in text
