"""Text extraction for AI assistant file uploads.

Supports PDF, DOCX, XLSX, PPTX, CSV. All extractors are synchronous and
should be called via `asyncio.to_thread(...)` from the route.

Text output is truncated to `MAX_CHARS` to protect the LLM context window
and MongoDB document size. The route is responsible for the size / MIME
gate — this module trusts its inputs.
"""
from __future__ import annotations

import csv
import io
import logging
import os
from typing import Tuple

logger = logging.getLogger(__name__)

# Cap extracted text so a single big spreadsheet can't blow the LLM context
# or the MongoDB doc limit. ~200k chars ≈ 50k tokens which is already
# aggressive; anything larger should be summarised by the assistant.
MAX_CHARS = 200_000

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".csv"}

EXT_TO_MIME = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".csv": "text/csv",
}


def _truncate(text: str) -> Tuple[str, bool]:
    if len(text) <= MAX_CHARS:
        return text, False
    return text[:MAX_CHARS] + "\n\n[... truncated by Zynthoro — file exceeds 200k characters ...]", True


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            logger.warning("pdf page %d extract failed: %s", i, e)
            continue
        page_text = page_text.strip()
        if page_text:
            parts.append(f"--- Page {i} ---\n{page_text}")
    return "\n\n".join(parts).strip()


def _extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    parts: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    # Include table cells too — a lot of business docs put tables in tables.
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    sheets: list[str] = []
    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in values):
                rows.append(",".join(values))
        if rows:
            sheets.append(f"--- Sheet: {ws.title} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets).strip()


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_parts.append(t)
            # Table shapes
            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
                except Exception:
                    pass
        if slide_parts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))
    return "\n\n".join(parts).strip()


def _extract_csv(content: bytes) -> str:
    # Try utf-8 first, fall back to latin-1 so we never crash on odd encodings.
    text = ""
    for enc in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            text = content.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if not text:
        return ""
    reader = csv.reader(io.StringIO(text))
    rows: list[str] = []
    for row in reader:
        if any((c or "").strip() for c in row):
            rows.append(",".join((c or "").strip() for c in row))
    return "\n".join(rows).strip()


_DISPATCH = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".csv": _extract_csv,
}


def extract_text(filename: str, content: bytes) -> Tuple[str, bool, str]:
    """Extract text from an uploaded file.

    Returns: (text, truncated_flag, mime_type).

    Raises ValueError for unsupported extensions.
    Raises RuntimeError if the underlying parser fails (route converts to HTTP 422).
    """
    _, ext = os.path.splitext(filename or "")
    ext = ext.lower()
    if ext not in _DISPATCH:
        raise ValueError(f"Unsupported file type: {ext or '(no extension)'}")
    try:
        text = _DISPATCH[ext](content)
    except Exception as e:  # noqa: BLE001
        logger.exception("File extraction failed for %s", filename)
        raise RuntimeError(f"Could not read {ext.upper()[1:]} file: {e}") from e
    text, truncated = _truncate(text)
    return text, truncated, EXT_TO_MIME[ext]
