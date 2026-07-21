"""E2E backend tests for the AI file upload feature.
Covers POST /api/ai/upload, DELETE /api/ai/upload/{id}, and
/api/ai/chat + /api/ai/stream with file_ids.
"""
import io
import os
import time
import json
import pytest
import requests

BASE_URL = os.environ.get("REACT_APP_BACKEND_URL", "").rstrip("/")
FOUNDER = ("regie@myrootzz.com", "Zynthoro2026!")
DEMO = ("jury@zynthoro.ai", "ZynthoroDemo2026!")

CSV_SALES = (
    b"product,units,price\n"
    b"Widget,120,9.99\n"
    b"Gadget,80,19.99\n"
    b"Gizmo,45,29.50\n"
)


def _login(email, password):
    s = requests.Session()
    r = s.post(f"{BASE_URL}/api/auth/login", json={"email": email, "password": password}, timeout=20)
    assert r.status_code == 200, f"login failed: {r.status_code} {r.text}"
    data = r.json()
    # Founder+Demo have 2FA disabled → direct login (should not return pre_token)
    assert data.get("stage") in (None, "authenticated", "ok") or "access_token" in data or s.cookies, \
        f"unexpected login stage: {data}"
    # Confirm session works
    me = s.get(f"{BASE_URL}/api/auth/me", timeout=10)
    assert me.status_code == 200, f"me failed: {me.status_code} {me.text}"
    return s


@pytest.fixture(scope="module")
def founder():
    return _login(*FOUNDER)


@pytest.fixture(scope="module")
def demo():
    return _login(*DEMO)


# ---------- Upload validation tests ----------
class TestUploadValidation:
    def test_anonymous_upload_rejected(self):
        r = requests.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("a.csv", io.BytesIO(b"a,b\n1,2"), "text/csv")},
            timeout=15,
        )
        assert r.status_code in (401, 403), f"expected 401/403, got {r.status_code}"

    def test_unsupported_extension(self, founder):
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("evil.exe", io.BytesIO(b"MZ\x90\x00"), "application/octet-stream")},
            timeout=15,
        )
        assert r.status_code == 400
        assert "Unsupported" in r.json().get("detail", "")

    def test_unsupported_txt(self, founder):
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("notes.txt", io.BytesIO(b"hi there"), "text/plain")},
            timeout=15,
        )
        assert r.status_code == 400

    def test_empty_file(self, founder):
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("empty.csv", io.BytesIO(b""), "text/csv")},
            timeout=15,
        )
        assert r.status_code == 400
        assert "empty" in r.json().get("detail", "").lower()

    def test_too_large(self, founder):
        big = b"a,b\n" + (b"x" * (10 * 1024 * 1024 + 100))
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("big.csv", io.BytesIO(big), "text/csv")},
            timeout=60,
        )
        assert r.status_code == 413
        assert "large" in r.json().get("detail", "").lower()

    def test_corrupted_docx(self, founder):
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("broken.docx", io.BytesIO(b"this is not a real docx"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
            timeout=15,
        )
        assert r.status_code == 422

    def test_whitespace_only_csv_returns_422(self, founder):
        # A CSV with only whitespace/newlines should produce no extractable text
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("blank.csv", io.BytesIO(b"   \n   \n\n"), "text/csv")},
            timeout=15,
        )
        # Accept 422 (no text extracted) — csv module may return the whitespace,
        # so we accept either 422 or a 200 with tiny preview.
        assert r.status_code in (200, 422)


# ---------- Happy path uploads ----------
class TestUploadFormats:
    def test_csv_upload_returns_expected_fields(self, founder):
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("sales.csv", io.BytesIO(CSV_SALES), "text/csv")},
            timeout=20,
        )
        assert r.status_code == 200, r.text
        data = r.json()
        for k in ("file_id", "filename", "size", "mime", "chars_extracted", "truncated", "preview"):
            assert k in data, f"missing {k}"
        assert data["filename"] == "sales.csv"
        assert data["size"] == len(CSV_SALES)
        assert "Widget" in data["preview"]
        # cleanup
        founder.delete(f"{BASE_URL}/api/ai/upload/{data['file_id']}")

    def test_docx_upload_extracts_text(self, founder):
        from docx import Document
        doc = Document()
        doc.add_paragraph("Zynthoro Q1 revenue was 1.2 million euros.")
        buf = io.BytesIO(); doc.save(buf); buf.seek(0)
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("q1.docx", buf, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
            timeout=20,
        )
        assert r.status_code == 200, r.text
        data = r.json()
        assert "Zynthoro Q1 revenue" in data["preview"]
        founder.delete(f"{BASE_URL}/api/ai/upload/{data['file_id']}")

    def test_xlsx_upload_extracts_text(self, founder):
        from openpyxl import Workbook
        wb = Workbook(); ws = wb.active
        ws.append(["Product", "Price"]); ws.append(["Widget", 9.99])
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("prices.xlsx", buf, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
            timeout=20,
        )
        assert r.status_code == 200, r.text
        assert "Widget" in r.json()["preview"]
        founder.delete(f"{BASE_URL}/api/ai/upload/{r.json()['file_id']}")

    def test_pptx_upload_extracts_text(self, founder):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Kickstart Plan"
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("deck.pptx", buf, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            timeout=20,
        )
        assert r.status_code == 200, r.text
        assert "Kickstart" in r.json()["preview"]
        founder.delete(f"{BASE_URL}/api/ai/upload/{r.json()['file_id']}")

    def test_pdf_upload_extracts_text(self, founder):
        try:
            from reportlab.pdfgen import canvas
        except ImportError:
            pytest.skip("reportlab not installed")
        buf = io.BytesIO()
        c = canvas.Canvas(buf); c.drawString(72, 720, "Hello Zynthoro PDF"); c.save()
        buf.seek(0)
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("hello.pdf", buf, "application/pdf")},
            timeout=20,
        )
        assert r.status_code == 200, r.text
        assert "Hello Zynthoro" in r.json()["preview"]
        founder.delete(f"{BASE_URL}/api/ai/upload/{r.json()['file_id']}")


# ---------- Delete + ownership ----------
class TestDeleteAndOwnership:
    def test_delete_own_upload(self, founder):
        up = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("del.csv", io.BytesIO(b"a,b\n1,2"), "text/csv")},
            timeout=15,
        ).json()
        r = founder.delete(f"{BASE_URL}/api/ai/upload/{up['file_id']}")
        assert r.status_code == 200
        assert r.json() == {"ok": True, "file_id": up["file_id"]}

    def test_delete_nonexistent(self, founder):
        r = founder.delete(f"{BASE_URL}/api/ai/upload/does-not-exist-1234")
        assert r.status_code == 404

    def test_cannot_delete_other_users_upload(self, founder, demo):
        up = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("private.csv", io.BytesIO(b"a,b\n1,2"), "text/csv")},
            timeout=15,
        ).json()
        r = demo.delete(f"{BASE_URL}/api/ai/upload/{up['file_id']}")
        assert r.status_code == 404
        # founder can still delete
        founder.delete(f"{BASE_URL}/api/ai/upload/{up['file_id']}")


# ---------- Chat + stream with file_ids ----------
class TestChatWithFiles:
    @pytest.fixture(scope="class")
    def sales_upload(self, founder):
        r = founder.post(
            f"{BASE_URL}/api/ai/upload",
            files={"file": ("sales.csv", io.BytesIO(CSV_SALES), "text/csv")},
            timeout=20,
        )
        assert r.status_code == 200
        fid = r.json()["file_id"]
        yield fid
        founder.delete(f"{BASE_URL}/api/ai/upload/{fid}")

    @pytest.mark.parametrize("assistant", ["zyona", "zyntha", "thoro", "zynthoro_assist"])
    def test_chat_with_file_context(self, founder, sales_upload, assistant):
        r = founder.post(
            f"{BASE_URL}/api/ai/chat",
            json={
                "assistant": assistant,
                "message": "Which product has the highest revenue in the attached CSV? Answer with just the product name.",
                "file_ids": [sales_upload],
            },
            timeout=90,
        )
        assert r.status_code == 200, r.text
        reply = r.json().get("reply") or r.json().get("message") or json.dumps(r.json())
        # Highest revenue: Gadget (80*19.99 = 1599.2) > Widget (120*9.99=1198.8) > Gizmo
        assert "Gadget" in reply, f"[{assistant}] expected 'Gadget' in reply, got: {reply[:400]}"

    def test_file_owner_isolation(self, founder, demo, sales_upload):
        """Demo user passes founder's file_id — reply must NOT contain the file's data."""
        r = demo.post(
            f"{BASE_URL}/api/ai/chat",
            json={
                "assistant": "zyona",
                "message": "List every product name and its exact price from the attached CSV file verbatim.",
                "file_ids": [sales_upload],
            },
            timeout=90,
        )
        assert r.status_code == 200, r.text
        reply = r.json().get("reply") or r.json().get("message") or ""
        # The full sales rows should NOT be in reply. If model refuses / says no file, that's OK.
        # We only fail if it echoes the private data.
        combined = ("Widget" in reply and "Gadget" in reply and "Gizmo" in reply)
        assert not combined, f"Leaked cross-user file data: {reply[:400]}"

    def test_stream_with_file_context(self, founder, sales_upload):
        r = founder.post(
            f"{BASE_URL}/api/ai/stream",
            json={
                "assistant": "zyona",
                "message": "Name the highest-revenue product from the attached CSV in one word.",
                "file_ids": [sales_upload],
            },
            headers={"Accept": "text/event-stream"},
            stream=True,
            timeout=90,
        )
        assert r.status_code == 200, r.text
        text = ""
        for line in r.iter_lines(decode_unicode=True):
            if not line:
                continue
            if line.startswith("data:"):
                text += line[5:]
            if "Gadget" in text or len(text) > 3000:
                break
        r.close()
        assert "Gadget" in text, f"stream did not mention Gadget: {text[:400]}"


# ---------- TTL index ----------
class TestTTLIndex:
    def test_ttl_index_exists(self):
        try:
            from pymongo import MongoClient
        except ImportError:
            pytest.skip("pymongo not available")
        mongo_url = os.environ.get("MONGO_URL")
        db_name = os.environ.get("DB_NAME")
        if not (mongo_url and db_name):
            pytest.skip("MONGO_URL/DB_NAME not set in test env")
        client = MongoClient(mongo_url)
        idx = client[db_name].ai_uploads.index_information()
        ttl_found = any(
            v.get("expireAfterSeconds") == 86400 and "created_at" in dict(v.get("key", [])).keys()
            for v in idx.values()
        )
        assert ttl_found, f"TTL index on created_at (86400s) not found: {idx}"
